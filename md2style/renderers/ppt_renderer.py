"""Ppt Renderer：标题驱动幻灯片生成器（数据接口执行层，第2层）。

基于 python-pptx；每个 # 标题生成一张新幻灯片（标题 + 内容）。
复用对应 html 风格（claude/mac）的主题色：背景、accent 强调色、标题条、卡片化内容、页码。
- 首个 H1 -> 封面页（大标题 + accent 下划线 + 副标题占位）
- 其余 H1/H2 -> 章节/内容页（accent 标题条 + 圆角内容卡片 + 阴影 + 页码）

样式取自 final_style 的 headings/body/accent_color/card_background/card_radius 等。
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .base import RendererBase, RendererDispatcher
from ..core.ir import IRNode, NodeType
from ..utils.color import hex_to_rgb


class PptRenderer(RendererBase):
    EMU = 914400

    def render(self, ir, final_style: dict, output_path: str) -> None:
        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        blank = prs.slide_layouts[6]

        headings = final_style.get("headings", {})
        body = final_style.get("body", {})
        accent = final_style.get("accent_color", "#888888")
        card_bg = final_style.get("card_background", "#F5F5F7")
        ar, ag, ab = hex_to_rgb(accent)
        bg_white = RGBColor(0xFF, 0xFF, 0xFF)

        slide = None
        bullets, code_blocks = [], []
        first_heading = True
        total = sum(1 for n in ir if n.type == NodeType.HEADING)
        idx = 0

        for node in ir:
            if node.type == NodeType.HEADING:
                if slide is not None:
                    self._fill(slide, bullets, code_blocks, body, accent, card_bg)
                idx += 1
                slide = prs.slides.add_slide(blank)
                self._paint_bg(slide, bg_white)
                if first_heading:
                    first_heading = False
                    self._cover(slide, node.text, accent)
                else:
                    self._title_bar(slide, node.text, headings, node.level, accent)
                bullets, code_blocks = [], []
            elif node.type in (NodeType.PARAGRAPH, NodeType.LIST, NodeType.QUOTE):
                bullets.append((node.text, node.meta.get("font")))
            elif node.type == NodeType.CODE:
                code_blocks.append(node.text)

        if slide is not None:
            self._fill(slide, bullets, code_blocks, body, accent, card_bg)
        # 给所有非封面页加页码
        self._add_page_numbers(prs, total)
        prs.save(output_path)

    # ---------- 页面构件 ----------
    def _paint_bg(self, slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _accent_bar(self, slide, accent):
        ar, ag, ab = hex_to_rgb(accent)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(0.75 * self.EMU)), Emu(int(0.67 * self.EMU)),
            Emu(int(3.2 * self.EMU)), Emu(int(0.09 * self.EMU)),
        )
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(ar, ag, ab)
        bar.line.fill.background()
        bar.shadow.inherit = False
        return bar

    def _cover(self, slide, title, accent):
        ar, ag, ab = hex_to_rgb(accent)
        # 居中大标题
        box = slide.shapes.add_textbox(Emu(int(1.2 * self.EMU)), Emu(int(2.4 * self.EMU)),
                                       Emu(int(9.8 * self.EMU)), Emu(int(1.6 * self.EMU)))
        tf = box.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = title
        run = p.runs[0]
        run.font.size = Pt(48); run.font.bold = True; run.font.name = "微软雅黑"
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        # 居中 accent 下划线
        ul = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(5.4 * self.EMU)), Emu(int(4.2 * self.EMU)),
            Emu(int(1.6 * self.EMU)), Emu(int(0.08 * self.EMU)),
        )
        ul.fill.solid(); ul.fill.fore_color.rgb = RGBColor(ar, ag, ab)
        ul.line.fill.background(); ul.shadow.inherit = False
        # 副标题占位
        sub = slide.shapes.add_textbox(Emu(int(2.0 * self.EMU)), Emu(int(4.5 * self.EMU)),
                                       Emu(int(8.0 * self.EMU)), Emu(int(0.8 * self.EMU)))
        stf = sub.text_frame; stf.word_wrap = True
        sp = stf.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
        sp.text = "md2style · " + accent
        sp.font.size = Pt(16); sp.font.name = "微软雅黑"
        sp.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    def _title_bar(self, slide, title, headings, level, accent):
        ar, ag, ab = hex_to_rgb(accent)
        self._accent_bar(slide, accent)
        box = slide.shapes.add_textbox(Emu(int(0.75 * self.EMU)), Emu(int(0.8 * self.EMU)),
                                       Emu(int(10.6 * self.EMU)), Emu(int(1.1 * self.EMU)))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title
        run = p.runs[0]
        hcfg = headings.get(f"h{level}", {})
        run.font.size = Pt(hcfg.get("size", 30))
        run.font.bold = True
        run.font.name = hcfg.get("font", "微软雅黑").split(",")[0].strip().strip("'\"")
        r, g, b = hex_to_rgb(hcfg.get("color", "#1A1A1A"))
        run.font.color.rgb = RGBColor(r, g, b)

    def _fill(self, slide, bullets, code_blocks, body, accent, card_bg):
        if not (bullets or code_blocks):
            return
        # 圆角内容卡片 + 阴影
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(0.75 * self.EMU)), Emu(int(2.05 * self.EMU)),
            Emu(int(10.6 * self.EMU)), Emu(int(4.4 * self.EMU)),
        )
        cr, cg, cb = hex_to_rgb(card_bg)
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(cr, cg, cb)
        card.line.fill.background()
        card.shadow.inherit = False
        try:
            card.shadow.type = 2  # 软阴影
        except Exception:
            pass
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(int(0.35 * self.EMU)); tf.margin_top = Emu(int(0.25 * self.EMU))
        first = True
        for b in bullets:
            text, font = b if isinstance(b, tuple) else (b, None)
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.text = "•  " + text
            self._style(p, body, font)
            p.space_after = Pt(10)
            first = False
        for code in code_blocks:
            p = tf.add_paragraph()
            p.text = code
            p.font.name = "Consolas"; p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_after = Pt(6)

    def _style(self, p, body, font=None):
        run = p.runs[0] if p.runs else None
        if run is None:
            return
        run.font.size = Pt(body.get("size", 16))
        run.font.name = (font or body.get("font", "微软雅黑")).split(",")[0].strip().strip("'\"")
        r, g, b = hex_to_rgb(body.get("color", "#1A1A1A"))
        run.font.color.rgb = RGBColor(r, g, b)

    def _add_page_numbers(self, prs, total):
        for i, slide in enumerate(prs.slides):
            # 跳过封面（第一页通常无标题条）
            if i == 0:
                continue
            box = slide.shapes.add_textbox(Emu(int(10.8 * self.EMU)), Emu(int(6.3 * self.EMU)),
                                           Emu(int(1.0 * self.EMU)), Emu(int(0.5 * self.EMU)))
            tf = box.text_frame
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
            p.text = f"{i + 1} / {total}"
            p.font.size = Pt(11); p.font.name = "微软雅黑"
            p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)


RendererDispatcher.register(".pptx", PptRenderer)
